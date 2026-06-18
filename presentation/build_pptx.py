"""
Build V2X platform showcase presentation as a .pptx file.
Run: python3 build_pptx.py
Output: V2X_Platform_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Colours ────────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0a, 0x0e, 0x1a)
BG_SLATE  = RGBColor(0x11, 0x18, 0x27)
ACCENT    = RGBColor(0x25, 0x63, 0xeb)
TEAL      = RGBColor(0x14, 0xb8, 0xa6)
AMBER     = RGBColor(0xf5, 0x9e, 0x0b)
RED       = RGBColor(0xef, 0x44, 0x44)
GREEN     = RGBColor(0x22, 0xc5, 0x5e)
WHITE     = RGBColor(0xf8, 0xfa, 0xfc)
LIGHT     = RGBColor(0xe2, 0xe8, 0xf0)
MUTED     = RGBColor(0x64, 0x74, 0x8b)
SURFACE   = RGBColor(0x1c, 0x28, 0x40)
BORDER    = RGBColor(0x1e, 0x3a, 0x5f)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completely blank


# ── Helpers ─────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(blank_layout)

def bg(slide, color=BG_DARK):
    """Fill slide background with a solid colour."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0.75)):
    """Add a rectangle shape."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        l, t, w, h
    )
    shape.line.fill.background()  # no line by default
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    return shape

def txt(slide, text, l, t, w, h,
        size=Pt(14), bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    """Add a text box."""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def para(tf, text, size=Pt(13), bold=False, color=LIGHT, align=PP_ALIGN.LEFT,
         italic=False, space_before=Pt(4)):
    """Add a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p

def accent_bar(slide, l, t, w=Inches(0.6), h=Pt(3), color=TEAL):
    b = box(slide, l, t, w, h, fill=color)
    return b

def eyebrow(slide, text, l=Inches(0.7), t=Inches(0.38)):
    txt(slide, text.upper(), l, t, Inches(9), Pt(16),
        size=Pt(11), bold=True, color=TEAL)

def heading(slide, text, l=Inches(0.7), t=Inches(0.62), w=Inches(11.9), size=Pt(36)):
    txt(slide, text, l, t, w, Inches(1.2),
        size=size, bold=True, color=WHITE)

def card(slide, l, t, w, h, fill=SURFACE, line_color=BORDER):
    return box(slide, l, t, w, h, fill=fill, line=line_color)

def bullet(slide, items, l, t, w, h,
           size=Pt(13), color=LIGHT, bullet_color=TEAL, spacing=Pt(6)):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = spacing
        run = p.add_run()
        run.text = "  •  " + item
        run.font.size  = size
        run.font.color.rgb = color
    return txb

def stat_block(slide, value, label, l, t, val_color=WHITE):
    txt(slide, value, l, t, Inches(2), Inches(0.7),
        size=Pt(32), bold=True, color=val_color, align=PP_ALIGN.CENTER)
    txt(slide, label.upper(), l, t + Inches(0.58), Inches(2), Inches(0.3),
        size=Pt(9), bold=True, color=MUTED, align=PP_ALIGN.CENTER)

def divider(slide, t, l=Inches(0.7), w=Inches(11.9)):
    box(slide, l, t, w, Pt(1), fill=BORDER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, BG_DARK)

# Subtle gradient-ish overlay rectangle
box(sl, Inches(0), Inches(0), Inches(6), SLIDE_H,
    fill=RGBColor(0x0d, 0x19, 0x2e))

# Teal accent bar
box(sl, Inches(0.7), Inches(1.9), Inches(0.08), Inches(2.6), fill=TEAL)

# Title
txt(sl, "V2X", Inches(0.9), Inches(1.9), Inches(10), Inches(1.1),
    size=Pt(72), bold=True, color=TEAL)
txt(sl, "Autonomous Emergency\nVehicle Priority Platform",
    Inches(0.9), Inches(2.85), Inches(9), Inches(1.6),
    size=Pt(36), bold=True, color=WHITE)

# Subtitle
txt(sl, "Authenticated V2X emergency signalling — from cryptographic session\n"
        "establishment to physical lane evasion — on real autonomous robots,\n"
        "on real hardware, with no simulation.",
    Inches(0.9), Inches(4.55), Inches(8.5), Inches(1.2),
    size=Pt(15), color=RGBColor(0x94, 0xa3, 0xb8), italic=True)

# Domain badges (as small coloured text boxes)
badge_defs = [
    ("Applied Cryptography", ACCENT,  RGBColor(0x06, 0x10, 0x27)),
    ("Computer Vision",      TEAL,    RGBColor(0x03, 0x1e, 0x1b)),
    ("Embedded Systems",     AMBER,   RGBColor(0x28, 0x1a, 0x03)),
    ("Autonomous Robotics",  GREEN,   RGBColor(0x05, 0x20, 0x0e)),
    ("V2X Networking",       RED,     RGBColor(0x27, 0x07, 0x07)),
]
bx = Inches(0.9)
for label, col, bg_col in badge_defs:
    box(sl, bx, Inches(5.95), Inches(1.9), Inches(0.32),
        fill=bg_col, line=col, line_w=Pt(0.5))
    txt(sl, label, bx + Inches(0.08), Inches(5.93), Inches(1.75), Inches(0.35),
        size=Pt(9), bold=True, color=col)
    bx += Inches(2.0)

txt(sl, "Siliris Technologies Pvt. Ltd.  ·  2026",
    Inches(0.9), Inches(6.8), Inches(8), Inches(0.4),
    size=Pt(10), color=MUTED)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE DEMONSTRATION
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
eyebrow(sl, "The Demo")
heading(sl, "What You See")
accent_bar(sl, Inches(0.7), Inches(1.55))

# Left column — steps
steps = [
    ("1", "An emergency vehicle approaches from behind",
     "It cryptographically authenticates itself with a roadside unit."),
    ("2", "Roadside unit issues a signed emergency event",
     "Encrypted broadcast sent to all vehicles on the subnet."),
    ("3", "Car cross-checks the ambulance's position",
     "Direct Pi-to-Pi UDP — no laptop needed. Decides autonomously to yield."),
    ("4", "Car moves to the inner edge",
     "Ambulance passes. Car recovers and resumes speed."),
]

ty = Inches(1.8)
for num, title, desc in steps:
    # Number circle
    box(sl, Inches(0.7), ty, Inches(0.38), Inches(0.38), fill=ACCENT)
    txt(sl, num, Inches(0.7), ty - Pt(1), Inches(0.38), Inches(0.38),
        size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, title, Inches(1.22), ty, Inches(5.4), Inches(0.28),
        size=Pt(13), bold=True, color=WHITE)
    txt(sl, desc, Inches(1.22), ty + Inches(0.27), Inches(5.4), Inches(0.28),
        size=Pt(11), color=RGBColor(0x94, 0xa3, 0xb8))
    ty += Inches(0.72)

# Right — track diagram card
card(sl, Inches(7.4), Inches(1.72), Inches(5.2), Inches(4.1),
     fill=RGBColor(0x06, 0x09, 0x0f))

track_lines = [
    "TRACK LAYOUT",
    "",
    "  ┌────────────────────────┐",
    "  │  outer yellow boundary │",
    "  │  ┌──────────────────┐  │",
    "  │  │  white oval lane │  │",
    "  │  │  ┌────────────┐  │  │",
    "  │  │  │yellow inner│  │  │",
    "  │  │  │island      │  │  │",
    "  │  │  │(yield zone)│  │  │",
    "  │  │  └────────────┘  │  │",
    "  │  └──────────────────┘  │",
    "  └────────────────────────┘",
    "",
    "  → Clockwise direction of travel",
    "  → 10 ft × 12 ft arena",
    "  → 10 AprilTag zones",
]
txb = slide_txb = sl.shapes.add_textbox(Inches(7.6), Inches(1.85), Inches(4.8), Inches(3.9))
tf = txb.text_frame
tf.word_wrap = False
first = True
for i, line in enumerate(track_lines):
    if first:
        p = tf.paragraphs[0]; first = False
    else:
        p = tf.add_paragraph()
    run = p.add_run()
    run.text = line
    if i == 0:
        run.font.size = Pt(9); run.font.bold = True; run.font.color.rgb = MUTED
    elif "yellow" in line.lower() and "boundary" not in line.lower() and "inner" in line.lower():
        run.font.size = Pt(10); run.font.color.rgb = AMBER
    elif "→" in line:
        run.font.size = Pt(10); run.font.color.rgb = GREEN
    else:
        run.font.size = Pt(10); run.font.color.rgb = RGBColor(0x7d, 0xd3, 0xfc)

# Key facts
card(sl, Inches(7.4), Inches(5.92), Inches(5.2), Inches(1.3),
     fill=RGBColor(0x06, 0x09, 0x18))
txt(sl, "No human presses a button at any stage   ·   No simulation   ·   9.4 ms auth latency",
    Inches(7.55), Inches(6.1), Inches(5.0), Inches(0.6),
    size=Pt(10), color=MUTED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — HARDWARE
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
eyebrow(sl, "Hardware")
heading(sl, "The Wolf Platform")
accent_bar(sl, Inches(0.7), Inches(1.55))

hw_cards = [
    ("🤖  Chassis", ACCENT,
     ["Powder-coated aluminium frame",
      "High-grip wheels",
      "Ring-lit camera mount",
      "Purpose-built — not a dev kit"]),
    ("⚡  Motor Controller", TEAL,
     ["Custom STM32F405 PCB (designed & fabricated)",
      "RTOS — guaranteed 20 Hz control loop",
      "Hardware encoders on all 4 wheels",
      "Lyra binary protocol over UART"]),
    ("📷  Sensing", AMBER,
     ["ArduCam 8MP IMX219",
      "Downward-facing at ~350 mm height",
      "Lane detection + AprilTag at 20 Hz",
      "320×240 px processing resolution"]),
    ("🖥️  Compute", GREEN,
     ["Raspberry Pi 5 on each robot",
      "Full autonomy stack on-device",
      "No cloud, no ROS at runtime",
      "C++ OBU + Python vision in parallel"]),
]

cx = Inches(0.65)
for title, col, items in hw_cards:
    card(sl, cx, Inches(1.72), Inches(3.0), Inches(5.35))
    # colour top bar
    box(sl, cx, Inches(1.72), Inches(3.0), Inches(0.06), fill=col)
    txt(sl, title, cx + Inches(0.15), Inches(1.85), Inches(2.7), Inches(0.38),
        size=Pt(13), bold=True, color=WHITE)
    ty = Inches(2.32)
    for item in items:
        txt(sl, "  •  " + item, cx + Inches(0.1), ty, Inches(2.8), Inches(0.33),
            size=Pt(11), color=RGBColor(0x94, 0xa3, 0xb8))
        ty += Inches(0.4)
    cx += Inches(3.18)

# Bottom stat row
divider(sl, Inches(7.2))
stats = [
    ("20 Hz", "Motor control loop", TEAL),
    ("4", "Hardware encoders", WHITE),
    ("RTOS", "Guaranteed timing", ACCENT),
    ("Custom PCB", "STM32F405 fabricated", AMBER),
    ("Pi 5", "On-device compute", GREEN),
]
sx = Inches(0.6)
for val, lbl, col in stats:
    txt(sl, val, sx, Inches(7.3), Inches(2.4), Inches(0.5),
        size=Pt(20), bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(sl, lbl.upper(), sx, Inches(7.72), Inches(2.4), Inches(0.28),
        size=Pt(8), color=MUTED, align=PP_ALIGN.CENTER)
    sx += Inches(2.5)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SYSTEM ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
eyebrow(sl, "Architecture")
heading(sl, "Two Systems. One Integration.")
accent_bar(sl, Inches(0.7), Inches(1.55))

# Three node boxes
nodes = [
    ("LAPTOP / DESKTOP PC", ACCENT,
     ["Desktop server (Flask + WebSocket)  :5000",
      "Issues cryptographic keys to all entities",
      "Runs the live authentication dashboard",
      "RSU binary (C++)  —  UDP :5000",
      "Authenticates OBU packets",
      "Broadcasts EMERGENCY_ACTIVE → :5001"]),
    ("CAR Pi  (v2x-car-01)", GREEN,
     ["Camera → lane follower → emergency handler → motors",
      "OBU subprocess (C++) — authenticates with RSU",
      "RSU alert listener  UDP :5001",
      "Position broadcaster  UDP :5002 ↔ ambulance",
      "Heartbeat → laptop every 15 s",
      "MJPEG vision stream  :5005"]),
    ("AMBULANCE Pi  (v2x-emgy)", RED,
     ["Camera → lane follower → drives straight through",
      "OBU re-authenticates every ~2 s",
      "Posts is_emergency: true → keeps ACTIVE alive",
      "Position broadcaster  UDP :5002 ↔ car",
      "Heartbeat → laptop every 15 s",
      "MJPEG vision stream  :5005"]),
]

cx = Inches(0.65)
for title, col, items in nodes:
    card(sl, cx, Inches(1.72), Inches(4.0), Inches(4.6))
    box(sl, cx, Inches(1.72), Inches(4.0), Inches(0.07), fill=col)
    txt(sl, title, cx + Inches(0.15), Inches(1.85), Inches(3.7), Inches(0.35),
        size=Pt(11), bold=True, color=col)
    ty = Inches(2.28)
    for item in items:
        txt(sl, "  ↳  " + item, cx + Inches(0.1), ty, Inches(3.8), Inches(0.32),
            size=Pt(10), color=LIGHT)
        ty += Inches(0.38)
    cx += Inches(4.22)

# Startup order strip
card(sl, Inches(0.65), Inches(6.45), Inches(12.0), Inches(0.75),
     fill=RGBColor(0x06, 0x09, 0x18))
txt(sl, "Startup order:   ① Desktop server   →   ② RSU binary   →   ③ Car Pi   →   ④ Ambulance Pi",
    Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.45),
    size=Pt(12), color=MUTED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — CRYPTOGRAPHIC AUTH
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
eyebrow(sl, "V2X Security")
heading(sl, "32-Step Mutual Authentication Protocol")
accent_bar(sl, Inches(0.7), Inches(1.55))

auth_steps = [
    ("1", TEAL,    "Registration (one-time)",
     "Desktop generates EC P-256 keypairs. Issues RID, AID, SK, PK_self, PK_RSU to each OBU."),
    ("2", ACCENT,  "AuthRequest  —  Ambulance OBU → RSU",
     "PID | timestamp | ECDH_encapsulate(PK_RSU) | ECDSA_sign(SK)"),
    ("3", ACCENT,  "RSU validates + derives session keys",
     "Timestamp ±500ms check → ECDSA verify → ECDH decapsulate → HKDF → SK_enc + SK_mac"),
    ("4", ACCENT,  "Mutual key confirmation KC1 / KC2",
     "Both sides prove they derived the same session keys — mutual authentication."),
    ("5", AMBER,   "Post-auth AES-256-GCM message",
     "{\"entity_id\": \"V2X_EMGY\", \"is_emergency\": true}  →  RSU broadcasts EMERGENCY_ACTIVE"),
]

ty = Inches(1.78)
for num, col, title, desc in auth_steps:
    box(sl, Inches(0.7), ty + Inches(0.04), Inches(0.36), Inches(0.36), fill=col)
    txt(sl, num, Inches(0.7), ty + Inches(0.02), Inches(0.36), Inches(0.36),
        size=Pt(11), bold=True, color=RGBColor(0x0a,0x0e,0x1a) if col==AMBER else WHITE,
        align=PP_ALIGN.CENTER)
    txt(sl, title, Inches(1.22), ty, Inches(6.0), Inches(0.28),
        size=Pt(12), bold=True, color=WHITE)
    txt(sl, desc, Inches(1.22), ty + Inches(0.28), Inches(6.2), Inches(0.28),
        size=Pt(10), color=RGBColor(0x94, 0xa3, 0xb8))
    ty += Inches(0.72)

# Right panel — why each mechanism
card(sl, Inches(7.6), Inches(1.72), Inches(5.05), Inches(4.0),
     fill=RGBColor(0x06, 0x09, 0x0f))
txt(sl, "WHY EACH MECHANISM EXISTS", Inches(7.75), Inches(1.85), Inches(4.7), Inches(0.3),
    size=Pt(9), bold=True, color=MUTED)

mechanisms = [
    ("Timestamp ±500ms",  "Prevents replay attacks"),
    ("ECDSA signature",   "Prevents packet forgery"),
    ("ECDH encapsulation","Prevents man-in-the-middle"),
    ("AES-256-GCM",       "Prevents eavesdropping"),
    ("HMAC-SHA-256",      "Prevents ciphertext bit-flip"),
    ("KC1/KC2 exchange",  "Ensures mutual key confirmation"),
]
ty = Inches(2.25)
for mech, reason in mechanisms:
    txt(sl, mech, Inches(7.75), ty, Inches(2.3), Inches(0.3),
        size=Pt(10), bold=True, color=ACCENT)
    txt(sl, reason, Inches(10.1), ty, Inches(2.4), Inches(0.3),
        size=Pt(10), color=LIGHT)
    ty += Inches(0.42)

# Stats
divider(sl, Inches(5.85))
stats2 = [("9.4 ms", "Auth latency", TEAL), ("32", "Protocol steps", WHITE), ("~2 s", "Re-auth cycle", AMBER)]
sx = Inches(0.7)
for val, lbl, col in stats2:
    txt(sl, val, sx, Inches(5.98), Inches(2.2), Inches(0.5),
        size=Pt(26), bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(sl, lbl.upper(), sx, Inches(6.43), Inches(2.2), Inches(0.28),
        size=Pt(8), color=MUTED, align=PP_ALIGN.CENTER)
    sx += Inches(2.3)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TWO SIGNALS
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
eyebrow(sl, "Emergency Logic")
heading(sl, "Two Independent Signals Must Both Be True")
accent_bar(sl, Inches(0.7), Inches(1.55))

# Signal A
card(sl, Inches(0.65), Inches(1.72), Inches(5.85), Inches(4.1))
box(sl, Inches(0.65), Inches(1.72), Inches(5.85), Inches(0.07), fill=ACCENT)
txt(sl, "🔐  SIGNAL A — V2X Authentication Chain", Inches(0.8), Inches(1.85),
    Inches(5.5), Inches(0.35), size=Pt(12), bold=True, color=ACCENT)
txt(sl, '"WHO is the emergency"', Inches(0.8), Inches(2.22),
    Inches(5.5), Inches(0.28), size=Pt(10), color=MUTED, italic=True)

sig_a = [
    "Ambulance OBU generates AuthRequest",
    "→  RSU validates (ECDSA + timestamp)",
    "→  ECDH session keys derived",
    "→  Post-auth: is_emergency: true",
    "→  RSU broadcasts EMERGENCY_ACTIVE",
    "→  UDP to 192.168.0.255:5001",
    "→  Car receives, sets _emergency = True",
]
ty = Inches(2.58)
for s in sig_a:
    c = ACCENT if s.startswith("→") else WHITE
    txt(sl, s, Inches(0.9), ty, Inches(5.3), Inches(0.3),
        size=Pt(10), color=c if not s.startswith("→  Car") else GREEN)
    ty += Inches(0.29)

# Signal B
card(sl, Inches(6.75), Inches(1.72), Inches(5.85), Inches(4.1))
box(sl, Inches(6.75), Inches(1.72), Inches(5.85), Inches(0.07), fill=TEAL)
txt(sl, "📡  SIGNAL B — Direct Position Sharing", Inches(6.9), Inches(1.85),
    Inches(5.5), Inches(0.35), size=Pt(12), bold=True, color=TEAL)
txt(sl, '"WHERE is the ambulance — no laptop needed"', Inches(6.9), Inches(2.22),
    Inches(5.5), Inches(0.28), size=Pt(10), color=MUTED, italic=True)

sig_b = [
    "Ambulance reads AprilTag zone",
    "→  JSON: {role: ambulance, zone: 7}",
    "→  UDP broadcast 192.168.0.255:5002",
    "→  10 Hz  —  purely Pi-to-Pi WiFi",
    "→  Car stores amb_zone = 7",
    "→  Freshness check: < 3 s old",
    "→  _position_known() = True",
]
ty = Inches(2.58)
for s in sig_b:
    c = TEAL if s.startswith("→") else WHITE
    txt(sl, s, Inches(6.9), ty, Inches(5.3), Inches(0.3),
        size=Pt(10), color=c if not s.startswith("→  Car") else GREEN)
    ty += Inches(0.29)

# Gate logic
card(sl, Inches(0.65), Inches(5.95), Inches(11.95), Inches(1.25),
     fill=RGBColor(0x06, 0x09, 0x10))
txt(sl, "YIELD DECISION:", Inches(0.85), Inches(6.05), Inches(1.8), Inches(0.3),
    size=Pt(10), bold=True, color=MUTED)
gate_items = [
    ("No V2X signal", "→ Follow lane normally", RED),
    ("V2X active, no position", "→ Wait for position fix", AMBER),
    ("Both active + amb ≤ 3 zones behind", "→ YIELD TRIGGERED", GREEN),
]
gx = Inches(2.0)
for cond, result, col in gate_items:
    txt(sl, cond, gx, Inches(6.05), Inches(3.5), Inches(0.28),
        size=Pt(10), bold=True, color=col)
    txt(sl, result, gx, Inches(6.35), Inches(3.5), Inches(0.28),
        size=Pt(10), color=LIGHT)
    gx += Inches(3.8)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — STATE MACHINE
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
eyebrow(sl, "Emergency Response")
heading(sl, "5-State Yield Sequence")
accent_bar(sl, Inches(0.7), Inches(1.55))

# State boxes with arrows
states = [
    ("NORMAL",     GREEN,  Inches(0.65)),
    ("EVADING",    AMBER,  Inches(3.15)),
    ("HOLDING",    RED,    Inches(5.65)),
    ("RECOVERING", ACCENT, Inches(8.15)),
    ("RESUMING",   TEAL,   Inches(10.65)),
]
STATE_FILLS = {
    "NORMAL":     RGBColor(0x04, 0x18, 0x0b),
    "EVADING":    RGBColor(0x1e, 0x14, 0x03),
    "HOLDING":    RGBColor(0x1e, 0x06, 0x06),
    "RECOVERING": RGBColor(0x06, 0x10, 0x27),
    "RESUMING":   RGBColor(0x03, 0x1e, 0x1b),
}
for name, col, sx in states:
    card(sl, sx, Inches(1.75), Inches(2.3), Inches(0.65),
         fill=STATE_FILLS[name], line_color=col)
    txt(sl, name, sx, Inches(1.75), Inches(2.3), Inches(0.65),
        size=Pt(13), bold=True, color=col, align=PP_ALIGN.CENTER)
    if sx != Inches(10.65):
        txt(sl, "→", sx + Inches(2.3), Inches(1.85), Inches(0.5), Inches(0.45),
            size=Pt(18), color=MUTED, align=PP_ALIGN.CENTER)

# Description cards
desc_cards = [
    ("EVADING", AMBER,
     "Turns right toward inner island (yellow line).\n"
     "Proportional yellow-repel steering.\n"
     "Exits when yellow line detected or 2 s timer fires."),
    ("HOLDING", RED,
     "Creeps slowly along inner island wall.\n"
     "Exits when: ambulance passes (amb_zone > own_zone)\n"
     "+ 1 s grace period  OR  V2X clears  OR  30 s timeout."),
    ("RECOVERING", ACCENT,
     "Arcs left back toward white line.\n"
     "Gyro mode: exits at 30° of measured yaw rotation.\n"
     "Timer fallback if IMU data is noisy."),
    ("RESUMING", TEAL,
     "Ramps speed 0 → normal over 2 s.\n"
     "Smooth re-entry to full-speed lane following.\n"
     "No jerk — gradual velocity ramp."),
]
cx = Inches(0.65)
for i, (title, col, desc) in enumerate(desc_cards):
    card(sl, cx, Inches(2.6), Inches(3.05), Inches(3.85))
    box(sl, cx, Inches(2.6), Inches(3.05), Inches(0.05), fill=col)
    txt(sl, title, cx + Inches(0.12), Inches(2.7), Inches(2.8), Inches(0.3),
        size=Pt(11), bold=True, color=col)
    txt(sl, desc, cx + Inches(0.12), Inches(3.1), Inches(2.82), Inches(3.0),
        size=Pt(10.5), color=LIGHT)
    cx += Inches(3.22)

# Key insight bar
card(sl, Inches(0.65), Inches(6.55), Inches(11.95), Inches(0.72),
     fill=RGBColor(0x06, 0x09, 0x10))
txt(sl, "Key: while HOLDING the car uses its frozen last-known zone + the ambulance's live UDP zone stream to "
        "detect when the ambulance has passed — no inner tags needed.",
    Inches(0.85), Inches(6.65), Inches(11.5), Inches(0.55),
    size=Pt(10), color=MUTED)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — COMPUTER VISION
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
eyebrow(sl, "Computer Vision")
heading(sl, "Real-Time Lane Following at 20 Hz")
accent_bar(sl, Inches(0.7), Inches(1.55))

# Left — algorithm steps
steps_cv = [
    ("①", "Grab 320×240 frame",         "Camera faces DOWN at the track surface"),
    ("②", "Crop top 40%",                "Remove far-field perspective distortion"),
    ("③", "HSV threshold — white line",  "Find centroid X of largest white contour"),
    ("④", "Proportional steering",        "wz = −kP × (centroid_x − centre) / (width/2)"),
    ("⑤", "Yellow repulsion",            "Override steering if yellow boundary too close"),
    ("⑥", "Detect AprilTag (every 3rd frame)", "Update zone → broadcast to peer via UDP"),
]
ty = Inches(1.78)
for num, title, desc in steps_cv:
    txt(sl, num, Inches(0.7), ty, Inches(0.4), Inches(0.28),
        size=Pt(12), bold=True, color=TEAL)
    txt(sl, title, Inches(1.18), ty, Inches(5.5), Inches(0.28),
        size=Pt(12), bold=True, color=WHITE)
    txt(sl, desc, Inches(1.18), ty + Inches(0.27), Inches(5.5), Inches(0.26),
        size=Pt(10), color=RGBColor(0x94, 0xa3, 0xb8))
    ty += Inches(0.66)

# Right — vision stream mockup
card(sl, Inches(7.4), Inches(1.72), Inches(5.25), Inches(3.6),
     fill=RGBColor(0x06, 0x09, 0x0f))
stream_lines = [
    "LIVE VISION STREAM  http://robot-ip:5005/",
    "─────────────────────────────────────────",
    " V2X_CAR_01                    23:14:05",
    " [Full camera frame 640px wide]",
    " ─── crop boundary ───────────────────",
    " [LANE overlay]  │  [HSV MASK]",
    "  green=centre   │   grey  = white px",
    "  red=centroid   │   yellow= yellow px",
    "──────────────────────────────────────",
    " CAR z=3  AMB z=7  vx=0.20  NORMAL",
    " ARMED  AUTO  V2X:ACTIVE  PI:52.3°C",
]
txb2 = sl.shapes.add_textbox(Inches(7.55), Inches(1.85), Inches(4.95), Inches(3.4))
tf2 = txb2.text_frame; tf2.word_wrap = False; first2 = True
for i, line in enumerate(stream_lines):
    if first2: p2 = tf2.paragraphs[0]; first2 = False
    else: p2 = tf2.add_paragraph()
    run2 = p2.add_run(); run2.text = line
    if i == 0: run2.font.size=Pt(8); run2.font.bold=True; run2.font.color.rgb=MUTED
    elif "ARMED" in line: run2.font.size=Pt(9); run2.font.color.rgb=GREEN; run2.font.bold=True
    elif "CAR z" in line: run2.font.size=Pt(9); run2.font.color.rgb=TEAL
    elif "─" in line: run2.font.size=Pt(8); run2.font.color.rgb=BORDER
    else: run2.font.size=Pt(9); run2.font.color.rgb=LIGHT

# Algorithms + failure modes
card(sl, Inches(7.4), Inches(5.45), Inches(5.25), Inches(1.75))
txt(sl, "3 SWAPPABLE ALGORITHMS", Inches(7.55), Inches(5.55), Inches(5.0), Inches(0.3),
    size=Pt(9), bold=True, color=MUTED)
algos = [("centroid", "Production default", GREEN),
         ("pure pursuit", "Geometric path", MUTED),
         ("recorded path", "Replay motor commands", MUTED)]
ax = Inches(7.55)
for name, desc2, col2 in algos:
    txt(sl, name, ax, Inches(5.88), Inches(1.6), Inches(0.3),
        size=Pt(11), bold=True, color=col2)
    txt(sl, desc2, ax, Inches(6.2), Inches(1.6), Inches(0.28),
        size=Pt(9), color=MUTED)
    ax += Inches(1.72)

# Failure modes
card(sl, Inches(0.65), Inches(6.5), Inches(6.55), Inches(0.75),
     fill=RGBColor(0x06, 0x09, 0x10))
txt(sl, "Failure modes are always safe:  "
        "no white >4s → slow  ·  no white >8s → stop  ·  no position → wait  ·  V2X drops → stop yielding",
    Inches(0.8), Inches(6.62), Inches(6.3), Inches(0.5),
    size=Pt(9), color=MUTED)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — DASHBOARD
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
eyebrow(sl, "Human-Machine Interface")
heading(sl, "Real-Time WebSocket Dashboard")
accent_bar(sl, Inches(0.7), Inches(1.55))

# Dashboard mockup
card(sl, Inches(0.65), Inches(1.72), Inches(6.0), Inches(4.8),
     fill=RGBColor(0x06, 0x09, 0x0f))
dash_lines = [
    "  V2X AUTHENTICATION DASHBOARD",
    "  ─────────────────────────────────────────────",
    "  ENTITIES: 3     (RSU · V2X · V2X_EMGY)",
    "",
    "  ●  RSU         ONLINE    —    No",
    "  ●  V2X         ONLINE    —    No",
    "  ●  V2X_EMGY   ONLINE    🚑   Yes",
    "",
    "  EMERGENCY PRIORITY GRANTS:   147",
    "  SESSION_ESTABLISHED:          12",
    "  POST_AUTH_HMAC_FAIL:           0",
    "  FAILURES:                      0",
    "  AVG AUTH LATENCY:           9.4 ms",
]
txb3 = sl.shapes.add_textbox(Inches(0.8), Inches(1.85), Inches(5.7), Inches(4.55))
tf3 = txb3.text_frame; tf3.word_wrap = False; first3=True
for i, line in enumerate(dash_lines):
    if first3: p3=tf3.paragraphs[0]; first3=False
    else: p3=tf3.add_paragraph()
    run3=p3.add_run(); run3.text=line
    if i==0: run3.font.size=Pt(10); run3.font.bold=True; run3.font.color.rgb=MUTED
    elif "ENTITIES" in line: run3.font.size=Pt(11); run3.font.bold=True; run3.font.color.rgb=GREEN
    elif "V2X_EMGY" in line: run3.font.size=Pt(10); run3.font.color.rgb=RED
    elif "●" in line: run3.font.size=Pt(10); run3.font.color.rgb=GREEN
    elif "GRANTS" in line: run3.font.size=Pt(10); run3.font.color.rgb=AMBER
    elif "LATENCY" in line: run3.font.size=Pt(10); run3.font.color.rgb=TEAL
    elif "─" in line: run3.font.size=Pt(8); run3.font.color.rgb=BORDER
    else: run3.font.size=Pt(10); run3.font.color.rgb=LIGHT

# What it shows
card(sl, Inches(6.9), Inches(1.72), Inches(5.7), Inches(4.8))
txt(sl, "WHAT THE DASHBOARD SHOWS", Inches(7.05), Inches(1.85), Inches(5.4), Inches(0.3),
    size=Pt(9), bold=True, color=MUTED)
dash_items = [
    "All entities online/offline with heartbeat watchdog",
    "Emergency flag per entity — live via WebSocket",
    "Session event counts — no page refresh needed",
    "Authentication latency — rolling average",
    "POST_AUTH_HMAC_FAIL — instant failure visibility",
    "Entity marks OFFLINE within ~50 s of crash",
    "SQLite audit log — full event history on disk",
]
ty2=Inches(2.28)
for item in dash_items:
    txt(sl, "  •  "+item, Inches(7.05), ty2, Inches(5.4), Inches(0.32),
        size=Pt(11), color=LIGHT)
    ty2+=Inches(0.4)

# Stats
divider(sl, Inches(6.6))
for val2, lbl2, col2, sx2 in [
    ("WebSocket", "Live push, no refresh", ACCENT, Inches(0.65)),
    ("10 Hz", "Position updates", TEAL, Inches(3.15)),
    ("SQLite", "Audit log", AMBER, Inches(5.65)),
    ("MJPEG", "Vision stream :5005", GREEN, Inches(8.15)),
    ("50 s", "Offline watchdog", RED, Inches(10.65)),
]:
    txt(sl, val2, sx2, Inches(6.72), Inches(2.3), Inches(0.45),
        size=Pt(18), bold=True, color=col2, align=PP_ALIGN.CENTER)
    txt(sl, lbl2.upper(), sx2, Inches(7.1), Inches(2.3), Inches(0.28),
        size=Pt(8), color=MUTED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — POST-QUANTUM
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
eyebrow(sl, "Future-Proofing")
heading(sl, "Post-Quantum Cryptography — Interface Ready")
accent_bar(sl, Inches(0.7), Inches(1.55))

txt(sl,
    "Classical ECDSA and ECDH are breakable by Shor's algorithm on a quantum computer.\n"
    "This platform abstracts all cryptography behind a single interface — switching to post-quantum\n"
    "requires changing one line in a config file. The 32-step protocol does not change.",
    Inches(0.7), Inches(1.78), Inches(7.8), Inches(0.85),
    size=Pt(12), color=RGBColor(0x94, 0xa3, 0xb8))

# Code card
card(sl, Inches(0.65), Inches(2.72), Inches(5.85), Inches(2.55),
     fill=RGBColor(0x06, 0x09, 0x0f))
code_lines = [
    ("// crypto_provider.h — pure virtual interface", MUTED),
    ("virtual encapsulate();", ACCENT),
    ("virtual decapsulate();", ACCENT),
    ("virtual sign();", ACCENT),
    ("virtual verify();", ACCENT),
    ("virtual derive_keys();", ACCENT),
    ("", WHITE),
    ('// rsu_config.json — one line to switch:', MUTED),
    ('"crypto_provider": "lattice"  // was: "placeholder"', GREEN),
]
txb4=sl.shapes.add_textbox(Inches(0.85),Inches(2.85),Inches(5.55),Inches(2.35))
tf4=txb4.text_frame; tf4.word_wrap=False; first4=True
for line4, col4 in code_lines:
    if first4: p4=tf4.paragraphs[0]; first4=False
    else: p4=tf4.add_paragraph()
    run4=p4.add_run(); run4.text=line4
    run4.font.size=Pt(10); run4.font.color.rgb=col4
    run4.font.name="Courier New"

# Planned algorithms
card(sl, Inches(0.65), Inches(5.35), Inches(5.85), Inches(1.85))
txt(sl, "PLANNED ALGORITHMS", Inches(0.85), Inches(5.48), Inches(5.5), Inches(0.28),
    size=Pt(9), bold=True, color=MUTED)
pq_items = [
    ("CRYSTALS-Dilithium", "Replaces ECDSA signatures  (NIST FIPS 204)"),
    ("CRYSTALS-Kyber",     "Replaces ECDH key encapsulation  (NIST FIPS 203)"),
    ("LWE-based hardness", "Resistant to Shor's algorithm on quantum computers"),
]
ty3=Inches(5.82)
for name3, desc3 in pq_items:
    txt(sl, name3, Inches(0.85), ty3, Inches(2.1), Inches(0.3),
        size=Pt(11), bold=True, color=TEAL)
    txt(sl, desc3, Inches(3.0), ty3, Inches(3.3), Inches(0.3),
        size=Pt(10), color=LIGHT)
    ty3+=Inches(0.38)

# Swappable layers
card(sl, Inches(6.75), Inches(2.72), Inches(5.85), Inches(4.5))
txt(sl, "ALL SWAPPABLE LAYERS", Inches(6.9), Inches(2.85), Inches(5.6), Inches(0.28),
    size=Pt(9), bold=True, color=MUTED)
layers = [
    ("Lane algorithm",    "centroid / pure pursuit / recorded-path — one line in config.yaml"),
    ("Motor controller",  "Lyra protocol over UART — retarget firmware without touching Python"),
    ("AprilTag layout",   "n_inner_tags / n_outer_tags — yield logic adapts automatically"),
    ("Crypto provider",   "placeholder (ECDH) → lattice (Kyber) — one JSON key"),
    ("Robot scale",       "Multiple cars + multiple ambulances — all work simultaneously"),
]
ty4=Inches(3.22)
for layer, ldesc in layers:
    txt(sl, layer, Inches(6.9), ty4, Inches(2.0), Inches(0.28),
        size=Pt(11), bold=True, color=TEAL)
    txt(sl, ldesc, Inches(8.95), ty4, Inches(3.5), Inches(0.28),
        size=Pt(10), color=LIGHT)
    ty4+=Inches(0.5)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — ENGINEERING SCOPE
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
eyebrow(sl, "Scope")
heading(sl, "Ten Engineering Domains. One System.")
accent_bar(sl, Inches(0.7), Inches(1.55))

domains = [
    ("Embedded Systems",    ACCENT,  "STM32F405, RTOS, hardware encoders × 4, custom-designed PCB"),
    ("Computer Vision",     TEAL,    "Real-time lane detection + AprilTag localisation under real lighting"),
    ("Robotics",            GREEN,   "5-state emergency evasion, closed-loop velocity, 3 swappable algorithms"),
    ("V2X Networking",      AMBER,   "OBU/RSU architecture, 32-step mutual auth, UDP broadcast"),
    ("Applied Cryptography",RED,     "ECDH, ECDSA, HKDF, AES-256-GCM, HMAC-SHA-256, PQ interface"),
    ("Distributed Systems", ACCENT,  "Multi-node coordination — no central controller, peer failure tolerant"),
    ("Real-Time Control",   TEAL,    "20 Hz RTOS motor loop, never preempted"),
    ("HMI",                 GREEN,   "WebSocket dashboard, MJPEG video stream with live telemetry overlay"),
    ("Edge Computing",      AMBER,   "Full stack on Raspberry Pi 5 — no cloud at runtime"),
    ("Autonomous Navigation",RED,    "Centroid, pure pursuit, recorded-path — tunable per deployment"),
]

ty5 = Inches(1.78)
col_break = 5
for i, (domain, col5, desc5) in enumerate(domains):
    lx = Inches(0.65) if i < col_break else Inches(6.75)
    ly = ty5 + Inches(0.62) * (i % col_break)
    box(sl, lx, ly + Inches(0.06), Inches(0.06), Inches(0.22), fill=col5)
    txt(sl, domain, lx + Inches(0.15), ly, Inches(2.1), Inches(0.28),
        size=Pt(11), bold=True, color=col5)
    txt(sl, desc5, lx + Inches(0.15), ly + Inches(0.27), Inches(5.8), Inches(0.28),
        size=Pt(10), color=LIGHT)

# Bottom callout
card(sl, Inches(0.65), Inches(6.52), Inches(11.95), Inches(0.72),
     fill=RGBColor(0x0d, 0x19, 0x2e), line_color=RGBColor(0x25, 0x63, 0xeb))
txt(sl,
    "This platform contains enough material for four independent final-year projects. "
    "This work required all four to exist simultaneously and operate as a single system.",
    Inches(0.85), Inches(6.62), Inches(11.5), Inches(0.55),
    size=Pt(11), color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — STATUS
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
eyebrow(sl, "Status")
heading(sl, "What's Live. What's Next.")
accent_bar(sl, Inches(0.7), Inches(1.55))

# Left — operational table
card(sl, Inches(0.65), Inches(1.72), Inches(5.85), Inches(5.3),
     fill=RGBColor(0x06, 0x09, 0x0f))
txt(sl, "OPERATIONAL NOW", Inches(0.85), Inches(1.88), Inches(5.5), Inches(0.28),
    size=Pt(9), bold=True, color=MUTED)

status_items = [
    ("V2X Auth — 32-step mutual protocol",  "Live",          GREEN),
    ("Robot mobility — lane following",      "Live",          GREEN),
    ("Emergency evasion — 5-state machine",  "Live",          GREEN),
    ("Real-time WebSocket dashboard",        "Live",          GREEN),
    ("Gyro-based RECOVERING exit",           "Live",          GREEN),
    ("Dead-reckoning position (encoders)",   "Live",          GREEN),
    ("End-to-end auth latency",              "9.4 ms",        TEAL),
    ("Post-quantum — LatticeProvider",       "Interface ready", AMBER),
]
ty6 = Inches(2.28)
for comp, status, col6 in status_items:
    txt(sl, comp, Inches(0.85), ty6, Inches(4.2), Inches(0.3),
        size=Pt(10.5), color=LIGHT)
    txt(sl, status, Inches(5.1), ty6, Inches(1.3), Inches(0.3),
        size=Pt(10), bold=True, color=col6, align=PP_ALIGN.RIGHT)
    ty6 += Inches(0.42)

# Right — next improvements
card(sl, Inches(6.75), Inches(1.72), Inches(5.85), Inches(5.3))
txt(sl, "FOUR TARGETED IMPROVEMENTS", Inches(6.9), Inches(1.88), Inches(5.6), Inches(0.28),
    size=Pt(9), bold=True, color=MUTED)

next_items = [
    ("1", ACCENT, "Geometry-based RECOVERING",
     "Integrate gyro_z from STM32 to measure actual degrees turned (30°) instead of waiting for a timer."),
    ("2", ACCENT, "Dead-reckoning between tags",
     "Continuous position from wheel encoder ticks; resets on every AprilTag sighting."),
    ("3", ACCENT, "Variable crop ratio by state",
     "NORMAL: 40%  ·  EVADING: 20%  ·  RECOVERING: 0%  — wider view during manoeuvres."),
    ("4", TEAL,   "Post-quantum — Kyber + Dilithium",
     "Plug into existing LatticeProvider interface. Protocol unchanged."),
]
ty7 = Inches(2.28)
for num7, col7, title7, desc7 in next_items:
    box(sl, Inches(6.9), ty7 + Inches(0.04), Inches(0.34), Inches(0.34), fill=col7)
    txt(sl, num7, Inches(6.9), ty7 + Inches(0.02), Inches(0.34), Inches(0.34),
        size=Pt(10), bold=True, color=RGBColor(0x0a,0x0e,0x1a), align=PP_ALIGN.CENTER)
    txt(sl, title7, Inches(7.35), ty7, Inches(5.1), Inches(0.3),
        size=Pt(11), bold=True, color=WHITE)
    txt(sl, desc7, Inches(7.35), ty7 + Inches(0.3), Inches(5.1), Inches(0.42),
        size=Pt(10), color=RGBColor(0x94, 0xa3, 0xb8))
    ty7 += Inches(0.95)

txt(sl, "No new hardware required for improvements 1–3. STM32 already sends gyro_z and wheel_ticks[4].",
    Inches(6.9), Inches(6.48), Inches(5.6), Inches(0.4),
    size=Pt(9), color=MUTED, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — CLOSING
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, RGBColor(0x06, 0x0c, 0x18))

# Background accent
box(sl, Inches(0), Inches(0), Inches(0.08), SLIDE_H, fill=TEAL)

txt(sl, "In one sentence", Inches(0.5), Inches(0.85), Inches(12.0), Inches(0.35),
    size=Pt(11), bold=True, color=TEAL, align=PP_ALIGN.CENTER)

txt(sl,
    "This project demonstrates authenticated V2X emergency\n"
    "vehicle priority — from cryptographic session establishment\n"
    "to physical lane evasion — on autonomous vehicles operating\n"
    "in the real world, on real hardware, with no simulation and\n"
    "no human intervention at any stage.",
    Inches(1.2), Inches(1.38), Inches(10.9), Inches(3.2),
    size=Pt(28), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

divider(sl, Inches(4.75), l=Inches(3.0), w=Inches(7.3))

txt(sl,
    "The cryptography is invisible to someone standing at the side of the table.\n"
    "What is visible is: the car moves, the ambulance passes, the car resumes.\n"
    "Nobody is thinking about the mathematics. They are thinking:  this works.",
    Inches(1.5), Inches(4.9), Inches(10.3), Inches(1.4),
    size=Pt(14), color=RGBColor(0x94, 0xa3, 0xb8), align=PP_ALIGN.CENTER, italic=True)

txt(sl, "Siliris Technologies Pvt. Ltd.  ·  2026",
    Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.35),
    size=Pt(10), color=MUTED, align=PP_ALIGN.CENTER)


# ── Save ────────────────────────────────────────────────────────────────────
out = "/workspace/presentation/V2X_Platform_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
